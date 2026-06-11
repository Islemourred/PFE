"""PPTX Engine for PFE Thesis Supervisor Meeting Presentation.
Academic style: Blue/Gold color scheme matching ESI SBA memoire."""

import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from PIL import Image as PILImage

FONT = "Calibri"
BLUE = RGBColor(0x29, 0x41, 0x7A)       # Thesis header blue
GOLD = RGBColor(0xC4, 0xA4, 0x4B)       # Thesis accent gold
LIGHT_BG = RGBColor(0xF5, 0xF6, 0xFA)   # Very light gray
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
DARK = RGBColor(0x2C, 0x2C, 0x2C)
GRAY = RGBColor(0x88, 0x88, 0x88)
LIGHT_BLUE = RGBColor(0xEA, 0xF0, 0xFB) # Light blue for tables
GREEN = RGBColor(0x2E, 0x8B, 0x57)

CHARTS_DIR = os.path.join(os.path.dirname(os.path.dirname(os.path.abspath(__file__))),
                           "output", "thesis_charts")
IMG_DIR = os.path.join(os.path.dirname(os.path.abspath(__file__)), "images")


class PFESlideBuilder:
    def __init__(self):
        self.prs = Presentation()
        self.prs.slide_width = Inches(13.33)
        self.prs.slide_height = Inches(7.5)
        self.num = 0

    def _bg(self, slide, color=WHITE):
        bg = slide.background
        fill = bg.fill
        fill.solid()
        fill.fore_color.rgb = color

    def _add_text(self, slide, left, top, width, height, text,
                  size=Pt(20), color=DARK, bold=False, align=PP_ALIGN.LEFT,
                  italic=False):
        tb = slide.shapes.add_textbox(Inches(left), Inches(top),
                                       Inches(width), Inches(height))
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = align
        run = p.add_run()
        run.text = text
        run.font.name = FONT
        run.font.size = size
        run.font.color.rgb = color
        run.font.bold = bold
        run.font.italic = italic
        return tb

    def _add_rich_text(self, slide, left, top, width, height, parts):
        """Add text with mixed formatting. parts = [(text, size, color, bold), ...]"""
        tb = slide.shapes.add_textbox(Inches(left), Inches(top),
                                       Inches(width), Inches(height))
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        for text, size, color, bold in parts:
            run = p.add_run()
            run.text = text
            run.font.name = FONT
            run.font.size = size
            run.font.color.rgb = color
            run.font.bold = bold
        return tb

    def _add_bullets(self, slide, left, top, width, height, bullets,
                     size=Pt(20), color=DARK, subs=None):
        from pptx.oxml.ns import qn
        from lxml import etree
        tb = slide.shapes.add_textbox(Inches(left), Inches(top),
                                       Inches(width), Inches(height))
        tf = tb.text_frame
        tf.word_wrap = True
        for i, b in enumerate(bullets):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.space_after = Pt(12)
            p.space_before = Pt(4)
            if b.strip():
                pPr = p._p.get_or_add_pPr()
                buNone = pPr.find(qn('a:buNone'))
                if buNone is not None:
                    pPr.remove(buNone)
                buClr = etree.SubElement(pPr, qn('a:buClr'))
                srgb = etree.SubElement(buClr, qn('a:srgbClr'))
                srgb.set('val', 'C4A44B')
                buChar = etree.SubElement(pPr, qn('a:buChar'))
                buChar.set('char', '\u25CF')
                buSz = etree.SubElement(pPr, qn('a:buSzPct'))
                buSz.set('val', '80000')
                pPr.set('marL', str(int(Pt(24))))
                pPr.set('indent', str(int(Pt(-18))))
            run = p.add_run()
            run.text = b
            run.font.name = FONT
            run.font.size = size
            run.font.color.rgb = color
            if subs and i in subs:
                for sb in subs[i]:
                    sp = tf.add_paragraph()
                    sp.space_after = Pt(6)
                    sp.space_before = Pt(2)
                    sp.level = 1
                    spPr = sp._p.get_or_add_pPr()
                    buClr2 = etree.SubElement(spPr, qn('a:buClr'))
                    srgb2 = etree.SubElement(buClr2, qn('a:srgbClr'))
                    srgb2.set('val', 'C4A44B')
                    buChar2 = etree.SubElement(spPr, qn('a:buChar'))
                    buChar2.set('char', '\u25B8')
                    spPr.set('marL', str(int(Pt(48))))
                    spPr.set('indent', str(int(Pt(-16))))
                    r2 = sp.add_run()
                    r2.text = sb
                    r2.font.name = FONT
                    r2.font.size = Pt(16)
                    r2.font.color.rgb = GRAY
        return tb

    def _title_bar(self, slide, title):
        line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                       Inches(0), Inches(0),
                                       Inches(13.33), Inches(0.06))
        line.fill.solid()
        line.fill.fore_color.rgb = GOLD
        line.line.fill.background()
        self._add_text(slide, 0.8, 0.3, 11, 0.7, title,
                       size=Pt(30), color=BLUE, bold=True)
        sep = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                      Inches(0.8), Inches(1.0),
                                      Inches(2.5), Inches(0.04))
        sep.fill.solid()
        sep.fill.fore_color.rgb = GOLD
        sep.line.fill.background()

    def _footer(self, slide):
        self.num += 1
        self._add_text(slide, 0.5, 7.0, 4, 0.3,
                       "PFE 2025/2026 | ESI Sidi Bel Abbes",
                       size=Pt(9), color=GRAY)
        self._add_text(slide, 11.5, 7.0, 1.5, 0.3,
                       str(self.num), size=Pt(9), color=GRAY,
                       align=PP_ALIGN.RIGHT)

    def _add_img(self, slide, img_path, left, top, max_w, max_h):
        if not img_path or not os.path.exists(img_path):
            return
        img = PILImage.open(img_path)
        iw, ih = img.size
        ratio = min(max_w / iw, max_h / ih)
        fw, fh = iw * ratio, ih * ratio
        cx = left + (max_w - fw) / 2
        cy = top + (max_h - fh) / 2
        slide.shapes.add_picture(img_path, Inches(cx), Inches(cy),
                                  Inches(fw), Inches(fh))

    # === SLIDE TYPES ===

    def title_slide(self, title, subtitle, authors, supervisors):
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self._bg(slide, WHITE)
        bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0,
                                      Inches(13.33), Inches(0.08))
        bar.fill.solid(); bar.fill.fore_color.rgb = GOLD; bar.line.fill.background()
        lb = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0,
                                     Inches(0.08), Inches(7.5))
        lb.fill.solid(); lb.fill.fore_color.rgb = BLUE; lb.line.fill.background()
        self._add_text(slide, 1.5, 1.5, 10, 1.2, title,
                       size=Pt(34), color=BLUE, bold=True, align=PP_ALIGN.CENTER)
        ol = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                     Inches(4.5), Inches(3.0),
                                     Inches(4.33), Inches(0.05))
        ol.fill.solid(); ol.fill.fore_color.rgb = GOLD; ol.line.fill.background()
        self._add_text(slide, 1.5, 3.3, 10, 0.7, subtitle,
                       size=Pt(20), color=DARK, align=PP_ALIGN.CENTER)
        self._add_text(slide, 1.5, 4.3, 10, 0.5, authors,
                       size=Pt(18), color=GRAY, align=PP_ALIGN.CENTER)
        self._add_text(slide, 1.5, 5.0, 10, 0.5, supervisors,
                       size=Pt(16), color=GRAY, align=PP_ALIGN.CENTER,
                       italic=True)
        self._add_text(slide, 1.5, 6.2, 10, 0.4,
                       "ESI Sidi Bel Abbes | Academic Year 2025/2026",
                       size=Pt(12), color=GRAY, align=PP_ALIGN.CENTER)
        self.num += 1

    def section_slide(self, section_num, title):
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self._bg(slide, BLUE)
        bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                      Inches(2), Inches(2.8),
                                      Inches(1.5), Inches(0.06))
        bar.fill.solid(); bar.fill.fore_color.rgb = GOLD; bar.line.fill.background()
        self._add_text(slide, 2, 1.8, 9, 0.8,
                       f"SECTION {section_num}",
                       size=Pt(16), color=GOLD, bold=True)
        self._add_text(slide, 2, 3.2, 9, 1.2, title,
                       size=Pt(38), color=WHITE, bold=True)
        self._footer(slide)

    def content(self, title, bullets, subs=None):
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self._bg(slide, WHITE)
        self._title_bar(slide, title)
        self._add_bullets(slide, 0.8, 1.3, 11.5, 5.5, bullets, subs=subs)
        self._footer(slide)
        return slide

    def split_text_img(self, title, bullets, img_path, subs=None):
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self._bg(slide, WHITE)
        self._title_bar(slide, title)
        self._add_bullets(slide, 0.8, 1.3, 6.0, 5.5, bullets, subs=subs)
        self._add_img(slide, img_path, 7.5, 1.3, 5.3, 5.3)
        self._footer(slide)
        return slide

    def chart_slide(self, title, chart_filename, caption=None):
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self._bg(slide, WHITE)
        self._title_bar(slide, title)
        chart_path = os.path.join(CHARTS_DIR, chart_filename)
        self._add_img(slide, chart_path, 0.8, 1.2, 11.5, 5.2)
        if caption:
            self._add_text(slide, 0.8, 6.6, 11.5, 0.4, caption,
                           size=Pt(11), color=GRAY, italic=True,
                           align=PP_ALIGN.CENTER)
        self._footer(slide)
        return slide

    def image_slide(self, title, img_path, caption=None):
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self._bg(slide, WHITE)
        self._title_bar(slide, title)
        self._add_img(slide, img_path, 0.8, 1.2, 11.5, 5.2)
        if caption:
            self._add_text(slide, 0.8, 6.6, 11.5, 0.4, caption,
                           size=Pt(11), color=GRAY, italic=True,
                           align=PP_ALIGN.CENTER)
        self._footer(slide)
        return slide

    def two_column(self, title, left_bullets, right_bullets,
                   left_title=None, right_title=None):
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self._bg(slide, WHITE)
        self._title_bar(slide, title)
        if left_title:
            self._add_text(slide, 0.8, 1.2, 5.5, 0.5, left_title,
                           size=Pt(18), color=BLUE, bold=True)
        if right_title:
            self._add_text(slide, 7.0, 1.2, 5.5, 0.5, right_title,
                           size=Pt(18), color=BLUE, bold=True)
        lt = 1.8 if left_title else 1.3
        self._add_bullets(slide, 0.8, lt, 5.5, 5.0, left_bullets, size=Pt(17))
        self._add_bullets(slide, 7.0, lt, 5.5, 5.0, right_bullets, size=Pt(17))
        sep = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                      Inches(6.5), Inches(1.3),
                                      Inches(0.03), Inches(5.0))
        sep.fill.solid(); sep.fill.fore_color.rgb = GOLD; sep.line.fill.background()
        self._footer(slide)
        return slide

    def end_slide(self, title, subtitle):
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self._bg(slide, BLUE)
        self._add_text(slide, 1.5, 2.5, 10, 1.2, title,
                       size=Pt(44), color=WHITE, bold=True,
                       align=PP_ALIGN.CENTER)
        bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                      Inches(5), Inches(4.0),
                                      Inches(3.33), Inches(0.05))
        bar.fill.solid(); bar.fill.fore_color.rgb = GOLD; bar.line.fill.background()
        self._add_text(slide, 1.5, 4.3, 10, 0.8, subtitle,
                       size=Pt(22), color=WHITE, align=PP_ALIGN.CENTER)
        self._footer(slide)

    def save(self, path):
        from pptx.oxml.ns import qn
        from lxml import etree
        for slide in self.prs.slides:
            sld = slide._element
            for existing in sld.findall(qn('p:transition')):
                sld.remove(existing)
            bg = slide.background
            is_section = False
            try:
                if bg.fill.fore_color.rgb == BLUE:
                    is_section = True
            except:
                pass
            transition = etree.SubElement(sld, qn('p:transition'))
            transition.set('spd', 'slow' if is_section else 'med')
            transition.set('advClick', '1')
            if is_section:
                push = etree.SubElement(transition, qn('p:push'))
            else:
                fade = etree.SubElement(transition, qn('p:fade'))
                fade.set('thruBlk', '0')
        self.prs.save(path)
        print(f"DONE: {path}")
        print(f"  Total slides: {self.num}")
